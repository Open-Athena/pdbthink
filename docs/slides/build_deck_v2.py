"""Build the pdbthink results deck from a collected results JSON.

Data-driven rather than hard-coded, because these numbers moved four times while
the sweep was running. Run ``collect_results.py`` first, then::

    python docs/slides/build_deck_v2.py results.json out.pptx
"""

from __future__ import annotations

import json
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

C = RGBColor.from_string
DEEP, SLATE, TEAL = "0B1A26", "15242E", "0E7C7B"
MINT, AMBER, CORAL = "4FD1C5", "E8A33D", "E2725B"
LIGHT, MUTED, LINE = "F7F9FA", "5B6E7A", "E1E8EC"
WHITE, PALE = "FFFFFF", "AFC3CE"
SERIF, SANS, MONO = "Cambria", "Calibri", "Consolas"

PRETTY = {
    "kimi_k3": "Kimi K3", "minimax_m3": "MiniMax M3", "qwen3_5_9b": "Qwen3.5 9B",
    "gpt_oss_120b": "gpt-oss-120b", "gpt_oss_20b": "gpt-oss-20b",
    "deepseek_v4_flash": "DeepSeek V4 Flash", "gemma_4_31b": "Gemma 4 31B",
    "marin": "Marin 32B (base)",
}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C(bg)
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, txt, x, y, w, h, *, size=15, bold=False, color=SLATE, font=SANS,
         align=PP_ALIGN.LEFT, italic=False, spacing=None, anchor=MSO_ANCHOR.TOP,
         space_after=0):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(txt if isinstance(txt, list) else [txt]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run(); run.text = line
        f = run.font
        f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
        f.color.rgb = C(color)
        if spacing is not None:
            f._rPr.set("spc", str(int(spacing * 100)))
    return box


def card(s, x, y, w, h, fill=WHITE, edge=LINE):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                           Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = C(fill)
    r.line.color.rgb = C(edge); r.line.width = Pt(0.75)
    r.shadow.inherit = False
    r.adjustments[0] = 0.04
    return r


def heading(s, title, kicker, color=SLATE, kicker_color=TEAL):
    text(s, kicker.upper(), 0.9, 0.62, 11.5, 0.3, size=11, bold=True,
         color=kicker_color, spacing=1.6)
    text(s, title, 0.9, 0.95, 11.5, 0.7, size=29, bold=True, color=color, font=SERIF)


def transparent(chart):
    """Charts default to a white plot area, which reads as a hole on a dark slide."""
    element = chart._chartSpace
    sp = element.find(qn("c:spPr"))
    if sp is None:
        sp = element.makeelement(qn("c:spPr"), {}); element.insert(0, sp)
    for child in list(sp):
        sp.remove(child)
    sp.append(sp.makeelement(qn("a:noFill"), {}))
    ln = sp.makeelement(qn("a:ln"), {})
    ln.append(sp.makeelement(qn("a:noFill"), {}))
    sp.append(ln)


def bar(s, cats, vals, x, y, w, h, *, colors=None, cat_color=SLATE, val_color=MUTED):
    data = CategoryChartData(); data.categories = cats; data.add_series("s", vals)
    chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y),
                               Inches(w), Inches(h), data).chart
    chart.has_title = False; chart.has_legend = False
    transparent(chart)
    plot = chart.plots[0]
    plot.gap_width = 55
    plot.has_data_labels = True
    lab = plot.data_labels
    lab.number_format = "0.00"; lab.number_format_is_linked = False
    lab.position = XL_LABEL_POSITION.OUTSIDE_END
    lab.font.size = Pt(11); lab.font.bold = True; lab.font.color.rgb = C(val_color)
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = C((colors or [TEAL] * len(cats))[i])
        point.format.line.fill.background()
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(11); ca.tick_labels.font.color.rgb = C(cat_color)
    ca.format.line.fill.background(); ca.major_tick_mark = XL_TICK_MARK.NONE
    va = chart.value_axis
    va.maximum_scale = 1.0; va.minimum_scale = 0
    va.has_major_gridlines = False; va.visible = False
    return chart


def main(results_path: str, out_path: str) -> None:
    data = json.loads(open(results_path).read())
    complete = [r for r in data["runs"] if r["complete"]]
    best = complete[0] if complete else data["runs"][0]
    label = PRETTY.get(best["label"], best["label"])

    # ------------------------------------------------------------- 1 title
    s = slide(DEEP)
    text(s, "pdbthink", 0.9, 2.15, 8, 0.9, size=52, bold=True, color=WHITE, font=SERIF)
    text(s, "Tool-free reasoning over protein structures", 0.9, 3.1, 9, 0.5,
         size=21, color=MINT)
    text(s, ["117 semantic instances  |  20 question families  |  49 protein groups",
             "59% of instances on structures released after 2023-01"],
         0.9, 3.95, 10, 0.9, size=14, color=PALE, space_after=6)
    text(s, "github.com/Open-Athena/pdbthink", 0.9, 6.5, 8, 0.3, size=12,
         color=TEAL, font=MONO)

    # ------------------------------------------------------- 2 the finding
    s = slide(DEEP)
    heading(s, "The scores are output budgets, not capabilities", "the finding",
            color=WHITE, kicker_color=MINT)
    text(s, [f"{label} scores {best['macro']:.3f} across the 20 families.",
             "",
             f"Restricted to responses that finished: {best['macro_finished']:.3f}."],
         0.9, 1.95, 6.5, 1.5, size=17, color=PALE, space_after=8)
    text(s, ["Truncation and inability score identically — zero.",
             "",
             "In every family where this model scores badly, the format-error",
             "count equals the truncation count. It is not failing to reason",
             "about interfaces; it never finished writing the answer.",
             "",
             "A benchmark of reasoning models that does not report truncation",
             "separately is reporting an output budget."],
         0.9, 3.7, 6.5, 2.7, size=13.5, color=PALE, space_after=3)
    cats = [PRETTY.get(r["label"], r["label"]) for r in complete][::-1]
    text(s, "AS SCORED", 7.85, 1.6, 4, 0.25, size=10.5, bold=True, color=MUTED, spacing=1.2)
    bar(s, cats, [r["macro"] for r in complete][::-1], 7.6, 1.85, 5.0, 2.1,
        colors=[TEAL] * len(cats), cat_color=PALE, val_color=PALE)
    text(s, "COMPLETED RESPONSES ONLY", 7.85, 4.2, 4.5, 0.25, size=10.5, bold=True,
         color=MINT, spacing=1.2)
    bar(s, cats, [r["macro_finished"] or 0 for r in complete][::-1], 7.6, 4.45, 5.0, 2.1,
        colors=[MINT] * len(cats), cat_color=PALE, val_color=MINT)

    # ---------------------------------------------------------- 3 all runs
    s = slide()
    heading(s, "Results", "all runs")
    y = 1.9
    for i, name in enumerate(["model", "renders", "macro", "finished only", "truncated"]):
        text(s, name.upper(), 1.15 + [0, 4.9, 6.5, 8.1, 10.3][i], y, 2.4, 0.3,
             size=10, bold=True, color=MUTED, spacing=1.2)
    y += 0.4
    for run in data["runs"]:
        full = run["complete"]
        card(s, 0.9, y - 0.07, 11.6, 0.5, fill=WHITE if full else "EEF2F4")
        name = PRETTY.get(run["label"], run["label"])
        if not full:
            name += f"   ({run['n_families']}/20 families — not comparable)"
        text(s, name, 1.15, y + 0.03, 4.8, 0.34, size=12.5, bold=full,
             color=SLATE if full else MUTED)
        text(s, str(run["n_renders"]), 6.05, y + 0.03, 1.4, 0.34, size=12.5,
             color=MUTED, font=MONO)
        text(s, f"{run['macro']:.3f}", 7.65, y + 0.03, 1.4, 0.34, size=12.5,
             bold=full, color=SLATE if full else MUTED, font=MONO)
        fin = run["macro_finished"]
        text(s, f"{fin:.3f}" if fin else "—", 9.25, y + 0.03, 1.4, 0.34,
             size=12.5, color=TEAL if full else MUTED, font=MONO)
        text(s, f"{run['truncated']}/{run['n_renders']}", 11.45, y + 0.03, 1.5, 0.34,
             size=12.5, color=CORAL if run["truncated"] else MUTED, font=MONO)
        y += 0.56
    text(s, "Greyed rows were cut short by a credit limit. Batches run in render-id "
            "order, so each covers a different alphabetically-early slice of the "
            "families — their macro averages are over different question sets and "
            "are not comparable to a complete run or to each other.",
         0.9, y + 0.2, 11.6, 0.8, size=11.5, italic=True, color=MUTED)

    # ----------------------------------------------------------- 4 the cliff
    s = slide()
    heading(s, f"{label}: solved outright, or never finished", "per family")
    per = best["per_family"]
    fams = sorted(per)
    solved = [f for f in fams if per[f]["score"] >= 0.95]
    text(s, f"{len(solved)} of {len(fams)} families at 0.95 or above",
         0.9, 1.85, 7, 0.35, size=15, bold=True, color=TEAL)
    text(s, "   ".join(solved), 0.9, 2.24, 7.2, 0.7, size=12.5, color=SLATE, font=MONO)
    weak = sorted([f for f in fams if per[f]["score"] < 0.95], key=lambda f: per[f]["score"])
    text(s, "the rest, and what they cost in truncation", 0.9, 3.1, 7, 0.3,
         size=13, bold=True, color=CORAL)
    y = 3.5
    for i, name in enumerate(["family", "score", "finished only", "truncated"]):
        text(s, name.upper(), 0.95 + [0, 1.6, 3.2, 5.4][i], y, 2, 0.3, size=10,
             bold=True, color=MUTED, spacing=1.2)
    y += 0.36
    for f in weak:
        v = per[f]
        text(s, f, 0.95, y, 1.4, 0.3, size=13, bold=True, color=SLATE, font=MONO)
        text(s, f"{v['score']:.3f}", 2.55, y, 1.4, 0.3, size=13, color=SLATE, font=MONO)
        fin = v["score_finished"]
        text(s, f"{fin:.3f}" if fin is not None else "—", 4.15, y, 1.4, 0.3,
             size=13, bold=True, color=TEAL, font=MONO)
        text(s, f"{v['truncated']}/{v['n']}", 6.35, y, 1.6, 0.3, size=13,
             color=CORAL, font=MONO)
        y += 0.38
    card(s, 8.4, 1.8, 4.1, 4.4)
    text(s, "Read the middle column", 8.75, 2.05, 3.6, 0.35, size=15, bold=True, color=SLATE)
    text(s, ["Every one of these families recovers once the",
             "responses that hit the cap are removed.",
             "",
             "The apparent weakness is the 32,768-token",
             "output budget, not the geometry.",
             "",
             "This is why truncation is its own failure",
             "category rather than folded into the score:",
             "a zero from a cut-off answer and a zero from",
             "a wrong one mean opposite things."],
         8.75, 2.55, 3.5, 3.4, size=12.5, color=MUTED, space_after=3)

    # ------------------------------------------------ 5 context-only control
    s = slide()
    heading(s, "What the coordinates are worth — and a trap", "controls")
    text(s, "Seven families and the six mechanistic episodes are also asked with the "
            "coordinates removed. For the automatic families the sanitised question "
            "names no protein, so what the control measures is the guessing floor.",
         0.9, 1.8, 11.5, 0.6, size=13.5, color=MUTED)
    base = best["context_only"]
    y = 2.6
    heads = ["family", "floor", "with coords", "naive gain", "truncated", "true gain"]
    for i, name in enumerate(heads):
        text(s, name.upper(), 0.95 + [0, 1.5, 2.9, 4.6, 6.4, 8.1][i], y, 2.2, 0.3,
             size=10, bold=True, color=MUTED, spacing=1.2)
    y += 0.36
    for f in sorted(base):
        v = base[f]
        text(s, f, 0.95, y, 1.2, 0.3, size=13, bold=True, color=SLATE, font=MONO)
        text(s, f"{v['floor']:.3f}", 2.45, y, 1.2, 0.3, size=13, color=SLATE, font=MONO)
        text(s, f"{v['with_coordinates']:.3f}", 3.85, y, 1.2, 0.3, size=13,
             color=SLATE, font=MONO)
        naive = v["gain_naive"]
        text(s, f"{naive:+.3f}", 5.55, y, 1.2, 0.3, size=13, font=MONO,
             bold=naive < 0, color=CORAL if naive <= 0 else MUTED)
        text(s, f"{100 * v['truncated_fraction']:.0f}%", 7.35, y, 1.2, 0.3,
             size=13, color=MUTED, font=MONO)
        cond = v["gain_conditioned"]
        text(s, f"{cond:+.3f}" if cond is not None else "—", 9.05, y, 1.4, 0.3,
             size=13, bold=True, color=TEAL, font=MONO)
        y += 0.38
    card(s, 0.9, y + 0.2, 11.6, 1.45, edge=CORAL)
    text(s, "The naive gain can invert its own sign", 1.2, y + 0.4, 8, 0.3,
         size=14, bold=True, color=CORAL)
    text(s, ["A context-only prompt is ~200 tokens and never truncates. A coordinate "
             "prompt runs to 87,500 tokens and truncates at up to 83%. Comparing them "
             "directly pits a finished cheap answer against a cut-off expensive one, and",
             "reports that coordinates make the model worse. Conditioned on completion, "
             "the same families show the gain the control was built to measure."],
         1.2, y + 0.75, 11.0, 0.8, size=12.5, color=MUTED, space_after=4)

    # ------------------------------------------------------------- 6 Marin
    s = slide(DEEP)
    heading(s, "Marin 32B: the benchmark does not fit in it", "a null result",
            color=WHITE, kicker_color=AMBER)
    text(s, ["Smallest prompt containing any coordinates", "Marin 32B context window"],
         0.9, 2.15, 5.6, 0.9, size=15, color=PALE, space_after=14)
    text(s, ["7,132 tokens", "4,096 tokens"], 6.6, 2.15, 3, 0.9, size=15,
         bold=True, color=MINT, font=MONO, space_after=14)
    text(s, "Not one structural question fits. All 48 renders that do are context-only "
            "controls, which by construction show no structure — so the run measures "
            "a guessing floor and nothing else.",
         0.9, 3.35, 11.5, 0.7, size=14, color=PALE)
    text(s, ["Two things make it unevaluable, both properties of the model:",
             "",
             "It is a base model. marin-32b-instruct does not exist. It answered P01 "
             "with the entire alphabet, and elsewhere repeated the prompt until it hit "
             "the cap — all 34 truncations.",
             "",
             "config.json claims rope_scaling to 65k; the model card says 4,096. vLLM "
             "overrides only behind a flag whose warning is that RoPE past the derived "
             "length returns NaN — and a zero from NaN is indistinguishable from a "
             "zero from inability."],
         0.9, 4.25, 11.5, 2.2, size=13, color=PALE, space_after=7)

    # ------------------------------------------------------------ 7 method
    s = slide()
    heading(s, "What makes these numbers reproducible", "method")
    items = [
        ("Gold answers are recomputed, never stored",
         "Every label comes from an oracle run on the coordinates the model was shown. A "
         "rotated variant, a cropped variant and a table variant must agree or the build fails."),
        ("The answers are not in the repository",
         "build regenerates them byte-identically and validate checks them against committed "
         "hashes. Two independent rebuilds were byte-identical."),
        ("Responses are content-addressed, not run-addressed",
         "Keyed on model, output budget and exact prompt text — not on render id. Adding a "
         "question costs the calls for that question. Full provider bodies are kept, so "
         "reasoning traces survive for inspection."),
        ("Truncation is a category, not a score",
         "Reported next to every number. An earlier model moved 0.475 → 0.710 with nothing "
         "changed but its output budget."),
    ]
    y = 1.9
    for title, body in items:
        card(s, 0.9, y, 11.6, 1.16)
        text(s, title, 1.25, y + 0.16, 10.9, 0.3, size=14.5, bold=True, color=SLATE)
        text(s, body, 1.25, y + 0.52, 10.9, 0.55, size=12.5, color=MUTED)
        y += 1.3

    # ----------------------------------------------------------- 8 caveats
    s = slide()
    heading(s, "What these numbers are not", "caveats")
    items = [
        ("Not curator-reviewed", "The 117-instance set is proposed, not accepted. A curator "
         "interface exists and decisions are being recorded; most of the set is pending."),
        ("One completion per prompt", "The protocol calls for three, and ten on a reliability "
         "subset. Intervals here cover benchmark composition only, not model stochasticity."),
        ("Three runs are incomplete", "A credit limit cut them off mid-batch. They cover "
         "different family subsets and are excluded from comparison rather than presented."),
        ("Budgets differ by model", "Each model got the largest output budget its context "
         "allows beside an 87,500-token prompt — gpt-oss 32,768 against Kimi's 65,536. "
         "Compare truncation counts, not only scores."),
        ("S02 is thin", "Two instances against a target of four: phosphorylation is rare in "
         "the pool and absent from all 49 FoldBench entries."),
    ]
    y = 1.9
    for title, body in items:
        text(s, title, 1.15, y, 3.5, 0.4, size=13.5, bold=True, color=CORAL)
        text(s, body, 4.9, y, 7.6, 0.7, size=12.5, color=MUTED)
        y += 0.94

    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
