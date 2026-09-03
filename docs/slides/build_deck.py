"""Build the pdbthink results deck."""

import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

C = RGBColor.from_string
DEEP, SLATE, TEAL = "0B1A26", "15242E", "0E7C7B"
MINT, AMBER, CORAL = "4FD1C5", "E8A33D", "E2725B"
LIGHT, MUTED, LINE = "F7F9FA", "5B6E7A", "E1E8EC"
WHITE, PALE, CARD_D = "FFFFFF", "AFC3CE", "12293A"

SERIF, SANS, MONO = "Cambria", "Calibri", "Courier New"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = C(bg)
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, txt, x, y, w, h, *, size=15, bold=False, color=SLATE, font=SANS,
         align=PP_ALIGN.LEFT, italic=False, spacing=None, anchor=MSO_ANCHOR.TOP, space_after=0):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = txt if isinstance(txt, list) else [txt]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        f = run.font
        f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
        f.color.rgb = C(color)
        if spacing is not None:
            f._rPr.set("spc", str(int(spacing * 100)))
    return box


def bullets(s, items, x, y, w, h, *, size=13, color=SLATE, gap=9):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = "•   " + item
        run.font.name, run.font.size = SANS, Pt(size)
        run.font.color.rgb = C(color)
    return box


def card(s, x, y, w, h, fill=WHITE, edge=LINE):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C(fill)
    shp.line.color.rgb = C(edge)
    shp.line.width = Pt(1)
    shp.adjustments[0] = 0.06
    shp.shadow.inherit = False
    return shp


def dot(s, x, y, d, color):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def heading(s, title, kicker, color=SLATE, kicker_color=TEAL):
    if kicker:
        text(s, kicker.upper(), 0.6, 0.38, 8, 0.25, size=11, bold=True,
             color=kicker_color, spacing=2)
    text(s, title, 0.6, 0.66, 12.1, 0.8, size=31, bold=True, color=color, font=SERIF)


def transparent(chart):
    """Drop the chart-area fill and border so it sits on the slide background."""
    from pptx.oxml.ns import qn
    cs = chart._chartSpace
    sp = cs.find(qn("c:spPr"))
    if sp is None:
        sp = cs.makeelement(qn("c:spPr"), {})
        plot_area = cs.find(qn("c:chart"))
        plot_area.addnext(sp)
    for child in list(sp):
        sp.remove(child)
    sp.append(sp.makeelement(qn("a:noFill"), {}))
    ln = sp.makeelement(qn("a:ln"), {})
    ln.append(ln.makeelement(qn("a:noFill"), {}))
    sp.append(ln)


def style_chart(chart, *, cat_color, val_color, grid, colors, labels=True,
                fmt="0.000", label_color=SLATE, cat_size=11, max_val=1.0):
    chart.has_legend = False
    chart.has_title = False
    transparent(chart)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.vary_by_categories = True
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = C(colors[i % len(colors)])
    if labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format, dl.number_format_is_linked = fmt, False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        dl.font.size, dl.font.name = Pt(12), SANS
        dl.font.color.rgb = C(label_color)
    ca, va = chart.category_axis, chart.value_axis
    for ax, col, size in ((ca, cat_color, cat_size), (va, val_color, 10)):
        ax.tick_labels.font.size = Pt(size)
        ax.tick_labels.font.name = SANS
        ax.tick_labels.font.color.rgb = C(col)
        ax.format.line.color.rgb = C(grid)
        ax.major_tick_mark = XL_TICK_MARK.NONE
        ax.minor_tick_mark = XL_TICK_MARK.NONE
    ca.has_major_gridlines = False
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = C(grid)
    va.major_gridlines.format.line.width = Pt(0.75)
    va.maximum_scale, va.minimum_scale = max_val, 0.0


def add_chart(s, kind, cats, vals, x, y, w, h, **kw):
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Macro score", vals)
    gf = s.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = gf.chart
    chart.font.name = SANS
    style_chart(chart, **kw)
    return chart


# ------------------------------------------------------------------ 1 title
s = slide(DEEP)
text(s, "pdbthink", 0.9, 1.95, 11, 1.0, size=54, bold=True, color=WHITE, font=SERIF)
text(s, "Tool-free reasoning over protein structures", 0.9, 2.95, 11, 0.6,
     size=24, color=MINT, font=SERIF)
text(s, "First results: a frontier reasoning model and a 1.7B local model on the smoke set",
     0.9, 3.68, 11, 0.5, size=16, color=PALE)
b = card(s, 0.9, 4.5, 3.25, 0.52, fill=TEAL, edge=TEAL)
text(s, "97 instances  ·  20 families", 0.9, 4.63, 3.25, 0.3, size=11, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)
text(s, "13 August 2026  ·  github.com/Open-Athena/pdbthink", 0.9, 6.45, 11, 0.35,
     size=11, color=MUTED)

# ------------------------------------------------------------ 2 what it tests
s = slide(LIGHT)
heading(s, "Six capabilities, measured separately", "what it tests")
rows = [
    ("Parsing", "Read a fixed-column coordinate file", "P01-P03", TEAL),
    ("Geometry", "Elementary 3D arithmetic on those numbers", "G01-G04", TEAL),
    ("Local structure", "Turn geometry into structural-biology concepts", "S01-S09", MINT),
    ("Interfaces & networks", "Reason across chains and contact graphs", "I01, N01", MINT),
    ("Two-state", "Compare two conformations", "T01", AMBER),
    ("Mechanism", "Connect a local change to a functional consequence", "6 episodes", AMBER),
]
y = 1.72
for name, what, fam, col in rows:
    dot(s, 0.62, y + 0.07, 0.3, col)
    text(s, name, 1.12, y, 3.3, 0.42, size=14.5, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, what, 4.5, y, 6.6, 0.42, size=13.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    text(s, fam, 11.2, y, 1.5, 0.42, size=12, color=col, font=MONO,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.73
text(s, "Each item shows one or two sanitised, randomly rotated structures and asks a question whose "
        "answer is recomputed from exactly those coordinates. Grading is deterministic throughout.",
     0.62, 6.35, 12.1, 0.6, size=13, italic=True, color=MUTED)

# ----------------------------------------------------------- 3 what was built
s = slide(LIGHT)
heading(s, "The candidate set", "what was built")
stats = [
    ("97", "semantic instances", "91 automatic + 6 mechanistic", TEAL),
    ("44", "protein groups", "from 52 PDB / AFDB entries", TEAL),
    ("183", "rendered variants", "rotation + representation controls", MINT),
    ("0", "validation errors", "across every consistency check", MINT),
]
x = 0.6
for value, label, sub, col in stats:
    card(s, x, 1.62, 2.95, 1.5)
    text(s, value, x, 1.76, 2.95, 0.7, size=34, bold=True, color=col, font=SERIF,
         align=PP_ALIGN.CENTER)
    text(s, label, x + 0.1, 2.48, 2.75, 0.28, size=11.5, bold=True, color=SLATE,
         align=PP_ALIGN.CENTER)
    text(s, sub, x + 0.1, 2.74, 2.75, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    x += 3.1
text(s, "Every gold answer is recomputed from the coordinates the model sees",
     0.6, 3.45, 12.1, 0.45, size=19, bold=True, font=SERIF)
bullets(s, [
    "Generators split into propose (find parameters that clear the ambiguity margins) and "
    "oracle (recompute the answer from any displayed structure)",
    "So a rotation variant, a cropped variant and a normalized-coordinate variant either agree, "
    "or the build fails",
    "3,270 rejection records across 25 distinct reasons — every near-tie and missing atom is "
    "logged, never silently resolved",
    "Rebuilds are byte-identical; answering every prompt with its own gold label scores exactly "
    "1.000 end to end",
], 0.62, 4.0, 12.1, 2.2, size=13.5)

# ---------------------------------------------------------------- 4 headline
s = slide(LIGHT)
heading(s, "DeepSeek V4 Flash vs a 1.7B local model", "headline")
add_chart(s, XL_CHART_TYPE.BAR_CLUSTERED,
          ["Qwen3-1.7B (local)", "DeepSeek V4 Flash"], [0.126, 0.710],
          0.6, 1.62, 7.1, 4.1,
          cat_color=SLATE, val_color=MUTED, grid=LINE, colors=[CORAL, TEAL], cat_size=12)
text(s, "Macro average across question families. 95% CI from a bootstrap clustered by protein.",
     0.6, 5.85, 7.1, 0.4, size=11, color=MUTED)
cards = [
    ("DeepSeek V4 Flash 0731", "0.710", "[0.603, 0.817]", "62 / 62 renders answered", TEAL),
    ("Qwen3-1.7B (local, ollama)", "0.126", "[0.057, 0.175]",
     "46 / 62 — rest exceed its 41k window", CORAL),
]
cy = 1.6
for name, score, ci, cov, col in cards:
    card(s, 8.0, cy, 4.72, 1.5)
    text(s, name, 8.22, cy + 0.14, 4.3, 0.3, size=11, bold=True, color=MUTED)
    text(s, score, 8.22, cy + 0.44, 1.6, 0.6, size=30, bold=True, color=col, font=SERIF)
    text(s, ci, 9.8, cy + 0.63, 2.7, 0.3, size=11, color=MUTED)
    text(s, cov, 8.22, cy + 1.06, 4.3, 0.3, size=10.5, color=SLATE)
    cy += 1.68
text(s, "11 of 20 families solved outright", 8.0, 5.05, 4.72, 0.35,
     size=14, bold=True, color=TEAL)
text(s, "exact accuracy 1.000", 8.0, 5.42, 4.72, 0.3, size=11, color=MUTED)
text(s, "Smoke set: 20 instances, 62 renders, one completion each.",
     0.6, 6.6, 12.1, 0.35, size=11, color=MUTED)

# -------------------------------------------------------------- 5 per family
s = slide(LIGHT)
heading(s, "Solved outright, or barely started", "where it succeeds and fails")
fams = ["G01", "G02", "G03", "N01", "P01", "P02", "P03", "S01", "S02", "S07", "S08",
        "MECH", "S05", "S09", "T01", "S03", "S04", "S06", "I01", "G04"]
vals = [1, 1, 1, 1, 1, 1, 1, 1, 1, 1, 1, .733, .667, .667, .419, .333, .333, .331, .239, 0]
cols = [TEAL] * 11 + [MINT] * 3 + [AMBER] * 1 + [CORAL] * 5
add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, fams, vals, 0.6, 1.6, 7.7, 4.35,
          cat_color=SLATE, val_color=MUTED, grid=LINE, colors=cols, labels=False, cat_size=9)
text(s, "DeepSeek V4 Flash, mean score per family", 0.6, 6.05, 7.7, 0.35, size=11, color=MUTED)
text(s, "The failures are consistent", 8.55, 1.62, 4.2, 0.4, size=18, bold=True, font=SERIF)
bullets(s, [
    "Set-valued questions — I01 interfaces, S06 ligand contacts, T01 — earn high partial "
    "credit and exact-set accuracy of 0.000. It finds most of the right residues and never quite "
    "the right set.",
    "G04 severe clash scores 0.000: it needs a global search over every atom pair.",
    "S03 burial and S04 secondary structure sit at 0.333 — these need SASA and DSSP-like "
    "reasoning, not distances.",
    "Mechanism episodes: 0.733 mean across fields, 0.000 exact. Some fields of every episode, "
    "all fields of none.",
], 8.55, 2.12, 4.2, 4.0, size=12.5, gap=11)

# ---------------------------------------------------------------- 6 controls
s = slide(LIGHT)
heading(s, "The controls behaved as designed", "validity checks")
blocks = [
    ("Coordinates are being read", "0.333 → 0.765", "context-only vs minimal PDB",
     "The mechanistic control strips the coordinates and keeps the setup. The gap is what the "
     "structure is worth — the model is not answering from recognition.", TEAL),
    ("No coordinate-frame sensitivity", "0.728 vs 0.804", "primary vs alternate rotation",
     "Distances and contacts are invariant by construction, so any gap would be the model rather "
     "than the task. There isn't one.", MINT),
    ("Format costs little here", "0.765 vs 0.696", "minimal PDB vs coordinate table",
     "Fixed-column PDB is not the bottleneck for a capable model — though it costs 41 tokens "
     "per atom against 22 for the table.", AMBER),
]
x = 0.6
for head, num, sub, body, col in blocks:
    card(s, x, 1.62, 3.95, 4.3)
    dot(s, x + 0.3, 1.92, 0.26, col)
    text(s, head, x + 0.68, 1.87, 3.0, 0.4, size=13.5, bold=True)
    text(s, num, x + 0.3, 2.42, 3.35, 0.7, size=26, bold=True, color=col, font=SERIF)
    text(s, sub, x + 0.3, 3.08, 3.35, 0.35, size=10.5, color=MUTED)
    text(s, body, x + 0.3, 3.55, 3.35, 2.2, size=12.5)
    x += 4.1
text(s, "Small n — one completion per render, 20 instances in 8 protein clusters. "
        "Directions, not effect sizes.", 0.6, 6.25, 12.1, 0.4, size=11, italic=True, color=MUTED)

# ----------------------------------------------------------------- 7 finding
s = slide(DEEP)
heading(s, "The score tripled without touching the model", "finding",
        color=WHITE, kicker_color=MINT)
add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED,
          ["8,192 tokens", "32,768 tokens", "65,536 tokens"], [0.475, 0.638, 0.710],
          0.85, 1.95, 6.5, 3.5,
          cat_color=PALE, val_color="5B7383", grid="1D3242",
          colors=[CORAL, AMBER, MINT], label_color=WHITE, max_val=0.8)
text(s, "output budget per response", 0.85, 5.55, 6.5, 0.3, size=11, color=MUTED,
     align=PP_ALIGN.CENTER)
text(s, "Truncation and inability score identically: zero.", 8.0, 2.0, 4.6, 0.9,
     size=19, bold=True, color=MINT, font=SERIF)
bullets(s, [
    "At 8k output tokens the model spent its whole budget on the reasoning trace and returned "
    "empty content for 32 of 62 renders.",
    "Truncations fell 32 → 17 → 7 as the budget rose. Seven remain at 64k, so 0.710 is "
    "still a lower bound.",
    "They concentrate in the set-valued families, where the trace grows with the structure.",
], 8.0, 3.0, 4.6, 2.6, size=13, color="D6E2E9", gap=11)
text(s, "Report the truncation count next to the score, or the score means nothing.",
     0.85, 6.5, 11.6, 0.4, size=14, italic=True, color=AMBER)

# ------------------------------------------------------------- 8 harness bugs
s = slide(LIGHT)
heading(s, "Three bugs that looked exactly like model failure", "harness, not model")
bugs = [
    ("Ollama dropped `think`",
     "Its OpenAI-compatible shim silently ignores the parameter. Qwen3 spent all 3,072 output "
     "tokens reasoning and returned empty content — 271 s per item, nothing to score.",
     "Added a native ollama_chat provider."),
    ("The parser read only one line",
     "The model wrote “Final Answer:” and put the value underneath. A correctly computed "
     "number scored as a format error — measuring layout, not reasoning.",
     "The marker line may now carry only a label."),
    ("Cloudflare blocked urllib",
     "Together sits behind a WAF that rejects the default Python-urllib agent. Every call failed "
     "403 on a valid key, which reads as bad credentials.",
     "Send a real User-Agent."),
]
y = 1.62
for i, (title, body, fix) in enumerate(bugs, 1):
    card(s, 0.6, y, 12.1, 1.48)
    dot(s, 0.9, y + 0.52, 0.42, CORAL)
    text(s, str(i), 0.9, y + 0.62, 0.42, 0.3, size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
    text(s, title, 1.52, y + 0.18, 3.7, 0.4, size=14.5, bold=True)
    text(s, body, 1.52, y + 0.58, 7.6, 0.82, size=12, color=MUTED)
    text(s, fix, 9.35, y + 0.52, 3.15, 0.6, size=12, bold=True, color=TEAL)
    y += 1.6
text(s, "All three move a score silently. They were catchable only because the runner records "
        "truncation, format errors and API errors as separate categories rather than folding them "
        "into “wrong”.", 0.6, 6.5, 12.1, 0.5, size=13, italic=True)

# ----------------------------------------------------------------- 9 caveats
s = slide(LIGHT)
heading(s, "What these numbers are not", "caveats")
items = [
    ("Not curator-reviewed", "The 97-instance set is proposed, not accepted. Three mechanistic "
     "episodes carry warnings where the published claim is not the top-ranked measurement.", AMBER),
    ("Not the full benchmark", "These runs use the 20-instance smoke set. The v1 candidate set is "
     "97 instances across 44 protein groups.", AMBER),
    ("Not a stable estimate", "One completion per render, 8 protein clusters. The protocol calls "
     "for three completions; run-to-run agreement is currently unmeasured.", CORAL),
    ("Not an upper bound", "Seven renders still truncate at a 64k output budget, each scored zero.",
     CORAL),
    ("Not contamination-controlled", "Release dates are recorded as a covariate. Most sources are "
     "decades-old structures.", MUTED),
]
y = 1.62
for head, body, col in items:
    dot(s, 0.62, y + 0.08, 0.28, col)
    text(s, head, 1.08, y, 3.2, 0.45, size=14, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, body, 4.35, y, 8.35, 0.85, size=13, color=MUTED)
    y += 1.0

# -------------------------------------------------------------------- 10 next
s = slide(DEEP)
heading(s, "To turn this into a result", "next", color=WHITE, kicker_color=MINT)
steps = [
    ("1", "Curate", "Review the 97 candidates and rule on the three episodes whose published claim "
     "the coordinates do not rank first."),
    ("2", "Run the full protocol", "Three completions per render on the accepted set, then the "
     "reliability subset, with agreement statistics."),
    ("3", "Widen the field", "Frontier models at several reasoning-effort levels, with the output "
     "budget set from the truncation count rather than guessed."),
]
y = 2.1
for n, head, body in steps:
    card(s, 0.85, y, 11.6, 1.2, fill=CARD_D, edge="1D3242")
    text(s, n, 1.15, y + 0.3, 0.5, 0.55, size=26, bold=True, color=MINT, font=SERIF)
    text(s, head, 1.85, y + 0.2, 3.0, 0.4, size=16, bold=True, color=WHITE)
    text(s, body, 1.85, y + 0.6, 10.2, 0.5, size=13, color=PALE)
    y += 1.4
text(s, "github.com/Open-Athena/pdbthink", 0.85, 6.5, 11.6, 0.4, size=14, color=MINT)

out = sys.argv[1] if len(sys.argv) > 1 else "pdbthink-results.pptx"
prs.save(out)
print("wrote", out, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
